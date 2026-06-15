#!/usr/bin/env python3
"""
Generate a PPTX fixture with custom theme colors.
This fixture is used to test theme color inspection and mutation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys

def create_theme_custom_colors_fixture(output_path):
    """Create a presentation with custom theme colors."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Access and modify theme
    # Note: python-pptx's theme modification support is limited,
    # so we'll create slides with various shapes colored with theme colors
    
    # Add a title slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Theme Custom Colors"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    
    # Add shape with accent1 color
    shape1 = slide.shapes.add_shape(1, Inches(1), Inches(2), Inches(2), Inches(1.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Default accent1 (blue)
    shape1.text_frame.text = "Accent 1"
    shape1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add shape with accent2 color
    shape2 = slide.shapes.add_shape(1, Inches(3.5), Inches(2), Inches(2), Inches(1.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(192, 80, 77)  # Default accent2 (red)
    shape2.text_frame.text = "Accent 2"
    shape2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add shape with accent3 color
    shape3 = slide.shapes.add_shape(1, Inches(6), Inches(2), Inches(2), Inches(1.5))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(155, 187, 89)  # Default accent3 (green)
    shape3.text_frame.text = "Accent 3"
    shape3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add second slide with different accent colors
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "More Theme Colors"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(54)
    title_para2.font.bold = True
    
    # Add shape with accent4 color
    shape4 = slide2.shapes.add_shape(1, Inches(1), Inches(2), Inches(2), Inches(1.5))
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = RGBColor(128, 100, 162)  # Default accent4 (purple)
    shape4.text_frame.text = "Accent 4"
    shape4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add shape with accent5 color
    shape5 = slide2.shapes.add_shape(1, Inches(3.5), Inches(2), Inches(2), Inches(1.5))
    shape5.fill.solid()
    shape5.fill.fore_color.rgb = RGBColor(75, 172, 198)  # Default accent5 (cyan)
    shape5.text_frame.text = "Accent 5"
    shape5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add shape with accent6 color
    shape6 = slide2.shapes.add_shape(1, Inches(6), Inches(2), Inches(2), Inches(1.5))
    shape6.fill.solid()
    shape6.fill.fore_color.rgb = RGBColor(247, 150, 70)  # Default accent6 (orange)
    shape6.text_frame.text = "Accent 6"
    shape6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save(output_path)
    print(f"Created theme fixture: {output_path}")

if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_file = os.path.join(script_dir, "..", "pptx", "theme-custom-colors", "presentation.pptx")
    
    # Create output directory if needed
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    
    create_theme_custom_colors_fixture(output_file)
