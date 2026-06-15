"""
Generates a PPTX with a slide containing high area occupancy (dense slide).
This fixture is used to test slide density metrics.
Output: testdata/pptx/layout-qa-dense-slide/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os


def generate_dense_slide():
    prs = Presentation()
    
    # Use a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Standard slide dimensions in inches: 10" x 7.5"
    slide_width = 10.0
    slide_height = 7.5
    
    # Fill most of the slide with shapes to create high density
    # Target: >70% area occupancy for "dense" classification
    
    # Add multiple rectangles that cover most of the slide
    colors = [
        RGBColor(255, 0, 0),     # Red
        RGBColor(0, 255, 0),     # Green
        RGBColor(0, 0, 255),     # Blue
        RGBColor(255, 255, 0),   # Yellow
        RGBColor(255, 0, 255),   # Magenta
        RGBColor(0, 255, 255),   # Cyan
    ]
    
    # Create a grid of shapes covering ~80% of the slide
    shape_width = Inches(3)   # 3 inches wide
    shape_height = Inches(2.4)  # 2.4 inches tall
    
    margin = Inches(0.2)  # Small margin
    
    color_index = 0
    x_pos = margin
    y_pos = margin
    
    while y_pos < Inches(slide_height - 0.5):
        x_pos = margin
        while x_pos < Inches(slide_width - 0.5):
            shape = slide.shapes.add_shape(
                1,  # Rectangle shape type
                x_pos, y_pos, shape_width, shape_height
            )
            shape.name = f"Shape_{color_index}"
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[color_index % len(colors)]
            shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
            
            x_pos += shape_width + Inches(0.1)
            color_index += 1
        
        y_pos += shape_height + Inches(0.1)
    
    # Create output directory
    output_dir = "testdata/pptx/layout-qa-dense-slide"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_dense_slide()
