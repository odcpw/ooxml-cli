"""
Generates a PPTX with a slide containing a picture placeholder.
Output: testdata/pptx/picture-placeholder/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches
import os


def generate_picture_placeholder():
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Picture Placeholder"
    
    # Add picture placeholder slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add picture placeholder shape
    left = Inches(1)
    top = Inches(1)
    width = Inches(4)
    height = Inches(3)
    pic_placeholder = slide2.shapes.add_picture(
        os.path.join(os.path.dirname(__file__), "../../test_image.png"),
        left, top, width, height
    ) if os.path.exists(os.path.join(os.path.dirname(__file__), "../../test_image.png")) else None
    
    if not pic_placeholder:
        # Create a simple test image if it doesn't exist
        try:
            from PIL import Image
            img = Image.new('RGB', (100, 100), color='red')
            img_path = os.path.join(os.path.dirname(__file__), "../../test_image.png")
            img.save(img_path)
            pic_placeholder = slide2.shapes.add_picture(img_path, left, top, width, height)
        except ImportError:
            # If PIL not available, just note the placeholder in a text box
            text_box = slide2.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = "[Picture Placeholder - requires image file]"
    
    # Create output directory
    output_dir = "testdata/pptx/picture-placeholder"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_picture_placeholder()
