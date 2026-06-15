#!/usr/bin/env python3
"""
Create rich text fixtures for testing text formatting, styles, and extraction.
"""
import os
import sys
import shutil
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree

def create_mixed_formatting_presentation():
    """Create a presentation with mixed bold, italic, and color formatting."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide with mixed formatting in subtitle
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Rich Text Formatting"
    
    # Add mixed formatting to subtitle
    tf = subtitle.text_frame
    tf.clear()  # Clear default text
    p = tf.paragraphs[0]
    
    # Normal text
    r = p.add_run()
    r.text = "This is "
    
    # Bold text
    r = p.add_run()
    r.text = "bold"
    r.font.bold = True
    
    # Normal
    r = p.add_run()
    r.text = " and this is "
    
    # Italic text
    r = p.add_run()
    r.text = "italic"
    r.font.italic = True
    
    # Normal
    r = p.add_run()
    r.text = " with "
    
    # Red colored text
    r = p.add_run()
    r.text = "color"
    r.font.color.rgb = RGBColor(255, 0, 0)
    
    # Normal
    r = p.add_run()
    r.text = "."
    
    # Slide 2: Multiple paragraphs with bullet levels
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "Bullet Points with Formatting"
    
    tf = content.text_frame
    tf.clear()
    
    # First bullet - level 0 with bold
    p = tf.paragraphs[0]
    p.level = 0
    r = p.add_run()
    r.text = "First item with "
    r = p.add_run()
    r.text = "bold formatting"
    r.font.bold = True
    
    # Second bullet - level 1 with color
    p = tf.add_paragraph()
    p.level = 1
    r = p.add_run()
    r.text = "Sub-item with "
    r = p.add_run()
    r.text = "blue color"
    r.font.color.rgb = RGBColor(0, 0, 255)
    
    # Third bullet - level 1 with italic
    p = tf.add_paragraph()
    p.level = 1
    r = p.add_run()
    r.text = "Another sub-item with "
    r = p.add_run()
    r.text = "italic"
    r.font.italic = True
    
    # Fourth bullet - level 0 with multiple colors
    p = tf.add_paragraph()
    p.level = 0
    r = p.add_run()
    r.text = "Second main item "
    r = p.add_run()
    r.text = "green"
    r.font.color.rgb = RGBColor(0, 128, 0)
    r = p.add_run()
    r.text = " and "
    r = p.add_run()
    r.text = "orange"
    r.font.color.rgb = RGBColor(255, 165, 0)
    
    return prs

def create_alignment_variants_presentation():
    """Create a presentation with different paragraph alignments and spacing."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title and alignment variants
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Paragraph Alignment Variants"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    
    # Left aligned
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Left aligned text"
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(18)
    
    # Center aligned
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Center aligned text"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    
    # Right aligned
    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Right aligned text"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(18)
    
    # Justify aligned
    left = Inches(0.5)
    top = Inches(4.5)
    width = Inches(9)
    height = Inches(2)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Justified text that should wrap across multiple lines to demonstrate the justify alignment option when we have enough content"
    p.alignment = PP_ALIGN.JUSTIFY
    p.font.size = Pt(14)
    
    # Slide 2: Font sizes and colors
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Font Sizes and Colors"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    
    # Different font sizes
    sizes = [12, 16, 20, 24, 28]
    colors = [
        RGBColor(0, 0, 0),      # Black
        RGBColor(255, 0, 0),    # Red
        RGBColor(0, 128, 0),    # Green
        RGBColor(0, 0, 255),    # Blue
        RGBColor(128, 0, 128),  # Purple
    ]
    
    for i, (size, color) in enumerate(zip(sizes, colors)):
        left = Inches(0.5)
        top = Inches(1.5 + i * 1.0)
        width = Inches(9)
        height = Inches(0.8)
        txBox = slide2.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Font size {size}pt"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    
    return prs

def create_bodypr_variants_presentation():
    """Create a presentation with different bodyPr (text body) properties."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Text Body Properties"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    
    # Different anchor/vertical align options
    top_positions = [1.5, 2.5, 3.5, 4.5, 5.5]
    anchors = [
        (MSO_ANCHOR.TOP, "Top anchor"),
        (MSO_ANCHOR.MIDDLE, "Middle anchor"),
        (MSO_ANCHOR.BOTTOM, "Bottom anchor"),
        (MSO_ANCHOR.TOP, "Top with word wrap"),
        (MSO_ANCHOR.MIDDLE, "Middle with word wrap"),
    ]
    
    for i, (anchor, label) in enumerate(anchors):
        left = Inches(0.5)
        top = Inches(top_positions[i])
        width = Inches(9)
        height = Inches(0.8)
        txBox = slide1.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.vertical_anchor = anchor
        if i >= 3:
            tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
    
    # Slide 2: Different margin/inset settings
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Margin Variants"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    
    # Different margin configurations
    margin_configs = [
        (Inches(0.2), Inches(0.2), Inches(0.2), Inches(0.2)),
        (Inches(0.5), Inches(0.1), Inches(0.5), Inches(0.1)),
        (Inches(0.1), Inches(0.5), Inches(0.1), Inches(0.5)),
    ]
    
    for i, (left_margin, top_margin, right_margin, bottom_margin) in enumerate(margin_configs):
        left = Inches(0.5)
        top = Inches(1.8 + i * 1.8)
        width = Inches(9)
        height = Inches(1.5)
        txBox = slide2.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.margin_left = left_margin
        tf.margin_top = top_margin
        tf.margin_right = right_margin
        tf.margin_bottom = bottom_margin
        p = tf.paragraphs[0]
        p.text = f"Margins: L={left_margin.inches:.1f}\" T={top_margin.inches:.1f}\" R={right_margin.inches:.1f}\" B={bottom_margin.inches:.1f}\""
        p.font.size = Pt(14)
    
    return prs

def create_numbered_list_presentation():
    """Create a presentation with numbered/bulleted lists with different formats."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Basic numbering
    bullet_slide_layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(bullet_slide_layout)
    title = slide1.shapes.title
    content = slide1.placeholders[1]
    
    title.text = "Numbered List"
    
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "First numbered item"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Second numbered item"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Sub-item of second"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Third numbered item"
    p.level = 0
    
    # Slide 2: Mixed bullet and numbering
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "Mixed Formatting"
    
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.level = 0
    r = p.add_run()
    r.text = "Item one with "
    r = p.add_run()
    r.text = "bold"
    r.font.bold = True
    
    p = tf.add_paragraph()
    p.level = 0
    r = p.add_run()
    r.text = "Item two with "
    r = p.add_run()
    r.text = "italic"
    r.font.italic = True
    r = p.add_run()
    r.text = " and "
    r = p.add_run()
    r.text = "red color"
    r.font.color.rgb = RGBColor(255, 0, 0)
    
    p = tf.add_paragraph()
    p.level = 1
    r = p.add_run()
    r.text = "Sub-item with "
    r = p.add_run()
    r.text = "blue"
    r.font.color.rgb = RGBColor(0, 0, 255)
    
    return prs

def save_presentation(prs, fixture_name, description):
    """Save a presentation to the testdata directory."""
    output_dir = Path("testdata/pptx") / fixture_name
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = output_dir / "presentation.pptx"
    prs.save(str(output_file))
    print(f"✓ Created {fixture_name}: {output_file}")
    print(f"  {description}")
    return output_file

def create_readme():
    """Create a README documenting the rich text fixtures."""
    readme_path = Path("testdata/pptx/README_RICH_TEXT.md")
    content = f"""# Rich Text Fixtures

This directory contains PPTX fixtures with rich text formatting for testing text extraction and formatting preservation.

Created: {datetime.now().isoformat()}

## Fixtures

### rich-formatting/presentation.pptx
- **Description**: Mixed formatting (bold, italic, colors) in title and content slides
- **Slides**: 2
- **Characteristics**:
  - Title slide with mixed formatting in subtitle (bold, italic, colored text)
  - Content slide with bullet points at different levels with various formatting
  - Tests extraction of run-level formatting properties (bold, italic, RGB color)

### rich-alignment/presentation.pptx
- **Description**: Different paragraph alignments and spacing variants
- **Slides**: 2
- **Characteristics**:
  - Slide 1: Left, center, right, and justify aligned paragraphs
  - Slide 2: Different font sizes (12pt to 28pt) and colors
  - Tests extraction of paragraph-level properties (alignment, font size, color)

### rich-bodypr/presentation.pptx
- **Description**: Different text body properties (anchoring, margins, word wrap)
- **Slides**: 2
- **Characteristics**:
  - Various vertical anchor settings (top, middle, bottom)
  - Different margin configurations (left, top, right, bottom)
  - Word wrap variants
  - Tests extraction of bodyPr-level properties (anchor, margins, word wrap)

### rich-numbered-lists/presentation.pptx
- **Description**: Numbered and mixed lists with multi-level formatting
- **Slides**: 2
- **Characteristics**:
  - Numbered list with multiple levels
  - Mixed bullet/number formatting with bold, italic, and colored text
  - Tests extraction of bullet/numbering properties and level hierarchy

## Testing Instructions

Test extract text command:

```bash
# Text output
./ooxml pptx extract text testdata/pptx/rich-formatting/presentation.pptx

# JSON output
./ooxml pptx extract text testdata/pptx/rich-formatting/presentation.pptx --format json

# With --rich flag (when implemented)
./ooxml pptx extract text testdata/pptx/rich-formatting/presentation.pptx --rich --format json

# Verify fixtures open in LibreOffice
libreoffice testdata/pptx/rich-formatting/presentation.pptx
libreoffice testdata/pptx/rich-alignment/presentation.pptx
libreoffice testdata/pptx/rich-bodypr/presentation.pptx
libreoffice testdata/pptx/rich-numbered-lists/presentation.pptx
```

## Property Extraction Coverage

### Run Properties (a:rPr)
- [x] Bold (a:rPr/@b)
- [x] Italic (a:rPr/@i)
- [x] Font size (a:rPr/@sz)
- [x] Font color (a:solidFill/a:srgbClr)
- [ ] Font name (a:rPr/@latin typeface)
- [ ] Underline (a:rPr/@u)
- [ ] Language (a:rPr/@lang)

### Paragraph Properties (a:pPr)
- [x] Alignment (a:pPr/@algn)
- [ ] Level (a:pPr/@lvl)
- [ ] Bullet character (a:buChar/@char)
- [ ] Numbering format (a:buAutoNum/@type)
- [ ] Spacing before/after (a:spcBef/a:spcAft)
- [ ] Line spacing (a:lnSpc)

### Body Properties (a:bodyPr)
- [x] Vertical anchor (a:bodyPr/@anchor)
- [x] Margin left (a:bodyPr/@lIns)
- [x] Margin top (a:bodyPr/@tIns)
- [x] Margin right (a:bodyPr/@rIns)
- [x] Margin bottom (a:bodyPr/@bIns)
- [x] Word wrap (a:bodyPr/@wrap)

## Running Tests

Generate/regenerate fixtures:
```bash
python3 testdata/generate/python/create_rich_text_fixtures.py
```

Run extract text tests:
```bash
go test ./internal/cli -run TestExtractText
go test ./internal/cli -run TestExtractTextRich
```

## Known Formatting Limitations

Some formatting options available in the PPTX format may not be fully captured by python-pptx:
- Some bullet/numbering options are limited
- Complex spacing configurations may vary
- Some color formats (theme colors, gradients) may not be available in python-pptx

See fixtures in LibreOffice Impress to verify actual formatting.
"""
    with open(readme_path, 'w') as f:
        f.write(content)
    print(f"✓ Created README: {readme_path}")
    return readme_path

def main():
    """Create all rich text fixtures."""
    print("Creating rich text fixtures...\n")
    
    # Create mixed formatting fixture
    prs = create_mixed_formatting_presentation()
    save_presentation(
        prs,
        "rich-formatting",
        "Mixed bold/italic/color text with bullet levels"
    )
    
    # Create alignment variants fixture
    prs = create_alignment_variants_presentation()
    save_presentation(
        prs,
        "rich-alignment",
        "Different paragraph alignments and font sizes/colors"
    )
    
    # Create bodyPr variants fixture
    prs = create_bodypr_variants_presentation()
    save_presentation(
        prs,
        "rich-bodypr",
        "Text body properties (anchor, margins, word wrap)"
    )
    
    # Create numbered list fixture
    prs = create_numbered_list_presentation()
    save_presentation(
        prs,
        "rich-numbered-lists",
        "Numbered and mixed lists with multi-level formatting"
    )
    
    # Create README
    create_readme()
    
    print("\n✓ All rich text fixtures created successfully!")
    print("\nNext steps:")
    print("1. Verify fixtures open in LibreOffice Impress")
    print("2. Extend text models with formatting properties (task-54)")
    print("3. Implement rich text extraction (task-55)")
    print("4. Add --rich flag to extract text CLI")
    print("5. Create golden JSON files for regression testing")
    print("\nRunning tests:")
    print("  go test ./internal/cli -run TestExtractText")

if __name__ == "__main__":
    main()
