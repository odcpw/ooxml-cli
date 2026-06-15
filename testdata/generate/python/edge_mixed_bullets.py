"""
Generates a PPTX with mixed bullet styles and levels (edge case).
Output: testdata/pptx/edge-mixed-bullets/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os


def generate_edge_mixed_bullets():
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Mixed Bullets Test"
    
    # Slide 2: Content with mixed bullets and levels
    content_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    title.text = "Mixed Bullet Levels"
    
    content = slide2.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Level 0: Regular bullet
    p = text_frame.paragraphs[0]
    p.text = "Level 0 - Main bullet"
    p.level = 0
    
    # Level 1: Sub-bullet
    p = text_frame.add_paragraph()
    p.text = "Level 1 - Sub-bullet"
    p.level = 1
    
    # Level 2: Sub-sub-bullet
    p = text_frame.add_paragraph()
    p.text = "Level 2 - Sub-sub-bullet"
    p.level = 2
    
    # Back to Level 0
    p = text_frame.add_paragraph()
    p.text = "Back to Level 0"
    p.level = 0
    
    # Level 1 again
    p = text_frame.add_paragraph()
    p.text = "Level 1 again"
    p.level = 1
    
    # Level 3: Deep nesting
    p = text_frame.add_paragraph()
    p.text = "Level 3 - Deep nesting"
    p.level = 3
    
    # Back to Level 0
    p = text_frame.add_paragraph()
    p.text = "Final Level 0"
    p.level = 0
    
    # Create output directory
    output_dir = "testdata/pptx/edge-mixed-bullets"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_edge_mixed_bullets()
