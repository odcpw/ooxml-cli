"""
Generates a PPTX with slides containing both speaker notes and handout content.
Output: testdata/pptx/notes-handout/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches
import os


def generate_notes_handout():
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Notes and Handout"
    
    # Add content slide with notes and handout
    content_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "Slide with Notes and Handout"
    
    text_frame = content.text_frame
    text_frame.text = "This slide has both speaker notes and is suitable for handout."
    
    # Add speaker notes
    notes_slide = slide2.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Speaker notes: Remember to mention the key points.\nAlso reference the handout for additional details."
    
    # Slide 3: Another slide with notes
    slide3 = prs.slides.add_slide(content_slide_layout)
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    title.text = "Second Slide with Notes"
    
    text_frame = content.text_frame
    text_frame.text = "Additional content for handout distribution"
    
    # Add notes to this slide
    notes_slide = slide3.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This slide provides additional information for attendees.\nInclude handout version with all details."
    
    # Create output directory
    output_dir = "testdata/pptx/notes-handout"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_notes_handout()
