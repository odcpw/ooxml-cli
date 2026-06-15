#!/usr/bin/env python3
"""
Generate a PPTX fixture with custom theme fonts.
This fixture is used to test theme font inspection and mutation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import sys

def create_theme_custom_fonts_fixture(output_path):
    """Create a presentation with custom theme fonts."""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add a title slide using a built-in layout
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Set title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Custom Theme Fonts"
    subtitle.text = "This presentation has custom fonts in the theme"
    
    # Add content slide
    content_slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide2 = prs.slides.add_slide(content_slide_layout)
    
    title2 = slide2.shapes.title
    title2.text = "Font Examples"
    
    # Add bullet points
    body_shape = slide2.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.text = "Major font: Used for titles"
    
    p = text_frame.add_paragraph()
    p.text = "Minor font: Used for body text"
    p.level = 0
    
    p = text_frame.add_paragraph()
    p.text = "Test different font sizes"
    p.level = 1
    
    # Add another slide with custom text
    slide3 = prs.slides.add_slide(content_slide_layout)
    title3 = slide3.shapes.title
    title3.text = "More Text"
    
    body_shape3 = slide3.placeholders[1]
    text_frame3 = body_shape3.text_frame
    text_frame3.text = "First point"
    
    p = text_frame3.add_paragraph()
    p.text = "Second point"
    
    p = text_frame3.add_paragraph()
    p.text = "Third point"
    
    # Save presentation
    prs.save(output_path)
    print(f"Created theme custom fonts fixture: {output_path}")

if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_file = os.path.join(script_dir, "..", "pptx", "theme-custom-fonts", "presentation.pptx")
    
    # Create output directory if needed
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    
    create_theme_custom_fonts_fixture(output_file)
