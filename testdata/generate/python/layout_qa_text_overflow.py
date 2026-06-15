"""
Generates a PPTX with a slide containing text that overflows its shape bounds.
This fixture is used to test text overflow detection heuristics.
Output: testdata/pptx/layout-qa-text-overflow/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import os


def generate_text_overflow():
    prs = Presentation()
    
    # Use a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add a small text box with large font and lots of text
    # This should trigger overflow detection
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(2)  # Small width
    height = Inches(1.5)  # Small height
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Add a paragraph with large font size and lots of text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(36)  # Large font
    p.text = "This is a very long text that will definitely not fit in this small text box because it has so much content and the font size is quite large making it impossible to display without overflow"
    
    # Add more paragraphs to ensure overflow
    for i in range(2):
        p = text_frame.add_paragraph()
        p.level = 0
        p.font.size = Pt(36)
        p.text = f"Additional paragraph {i+1} with even more text that definitely overflows"
    
    # Create output directory
    output_dir = "testdata/pptx/layout-qa-text-overflow"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_text_overflow()
