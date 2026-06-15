"""
Generates a PPTX with title slide and content layout slide.
Output: testdata/pptx/title-content/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Pt
import os


def generate_title_content():
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Title Content Presentation"
    subtitle.text = "Subtitle goes here"
    
    # Slide 2: Content slide
    content_slide_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "Content Slide"
    
    text_frame = content.text_frame
    text_frame.text = "This is the main content area"
    
    # Create output directory
    output_dir = "testdata/pptx/title-content"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_title_content()
