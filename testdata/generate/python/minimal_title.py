"""
Generates a minimal valid PPTX with a single title slide.
Output: testdata/pptx/minimal-title/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Pt
import os


def generate_minimal_title():
    prs = Presentation()
    
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = "Minimal Title Slide"
    
    # Create output directory
    output_dir = "testdata/pptx/minimal-title"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_minimal_title()
