"""
Generates a PPTX with a slide containing a table.
Output: testdata/pptx/table-slide/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os


def generate_table_slide():
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Table Slide"
    
    # Add table slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add table
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Fill table with data
    for i in range(rows):
        for j in range(cols):
            cell = table_shape.cell(i, j)
            cell.text = f"R{i}C{j}"
    
    # Create output directory
    output_dir = "testdata/pptx/table-slide"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_table_slide()
