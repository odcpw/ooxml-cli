"""
Generates a PPTX with slides using at least two distinct layout names.
Output: testdata/pptx/multi-layout/presentation.pptx
"""

from pptx import Presentation
import os


def generate_multi_layout():
    prs = Presentation()
    
    # Slide 1: Title Slide layout
    title_slide_layout = prs.slide_layouts[0]  # "Title Slide"
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Multi-Layout Presentation"
    
    # Slide 2: Title and Content layout
    content_slide_layout = prs.slide_layouts[1]  # "Title and Content"
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "Content Slide"
    text_frame = content.text_frame
    text_frame.text = "This slide uses Title and Content layout"
    
    # Slide 3: Section Header layout (different from Title Slide and Title and Content)
    if len(prs.slide_layouts) > 2:
        section_layout = prs.slide_layouts[2]  # "Section Header"
    else:
        section_layout = prs.slide_layouts[1]
    
    slide3 = prs.slides.add_slide(section_layout)
    title = slide3.shapes.title
    title.text = "Section Header"
    
    # Slide 4: Blank layout
    blank_layout = prs.slide_layouts[6]  # "Blank"
    slide4 = prs.slides.add_slide(blank_layout)
    
    # Add text to blank slide manually
    left = prs.slide_width // 4
    top = prs.slide_height // 4
    width = prs.slide_width // 2
    height = prs.slide_height // 2
    text_box = slide4.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "Blank Layout"
    
    # Create output directory
    output_dir = "testdata/pptx/multi-layout"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")
    print(f"  Layout 1: {prs.slide_layouts[0].name}")
    print(f"  Layout 2: {prs.slide_layouts[1].name}")
    if len(prs.slide_layouts) > 2:
        print(f"  Layout 3: {prs.slide_layouts[2].name}")


if __name__ == "__main__":
    generate_multi_layout()
