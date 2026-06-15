#!/usr/bin/env python3
"""
Create PPTX fixtures with specific geometry properties for testing.
Generates deterministic fixtures with known rotation, flip, and crop values.
"""
import os
import sys
import shutil
import zipfile
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

# XML Namespaces
NSMAP = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def modify_pptx_with_rotation(pptx_path, rotation_degrees):
    """
    Modify a PPTX file to add rotation to the first image.
    Rotation is in 1/60000 of a degree.
    """
    rotation_emu = rotation_degrees * 60000
    
    temp_dir = tempfile.mkdtemp()
    try:
        # Extract PPTX
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Modify slide1.xml to add rotation to first picture
        slide_path = os.path.join(temp_dir, 'ppt/slides/slide1.xml')
        if os.path.exists(slide_path):
            tree = etree.parse(slide_path)
            root = tree.getroot()
            
            # Find first picture
            pics = root.xpath('.//p:pic', namespaces=NSMAP)
            if pics:
                pic = pics[0]
                spPr = pic.find('.//p:spPr', NSMAP)
                if spPr is not None:
                    xfrm = spPr.find('.//a:xfrm', NSMAP)
                    if xfrm is not None:
                        xfrm.set('rot', str(rotation_emu))
            
            # Write back
            tree.write(slide_path, xml_declaration=True, encoding='UTF-8', standalone=True)
        
        # Repack PPTX
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root_dir, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_ref.write(file_path, arcname)
    finally:
        shutil.rmtree(temp_dir)

def modify_pptx_with_flip(pptx_path, flip_h=False, flip_v=False):
    """
    Modify a PPTX file to add flip properties to the first image.
    """
    temp_dir = tempfile.mkdtemp()
    try:
        # Extract PPTX
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Modify slide1.xml to add flip to first picture
        slide_path = os.path.join(temp_dir, 'ppt/slides/slide1.xml')
        if os.path.exists(slide_path):
            tree = etree.parse(slide_path)
            root = tree.getroot()
            
            # Find first picture
            pics = root.xpath('.//p:pic', namespaces=NSMAP)
            if pics:
                pic = pics[0]
                spPr = pic.find('.//p:spPr', NSMAP)
                if spPr is not None:
                    xfrm = spPr.find('.//a:xfrm', NSMAP)
                    if xfrm is not None:
                        if flip_h:
                            xfrm.set('flipH', '1')
                        if flip_v:
                            xfrm.set('flipV', '1')
            
            # Write back
            tree.write(slide_path, xml_declaration=True, encoding='UTF-8', standalone=True)
        
        # Repack PPTX
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root_dir, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_ref.write(file_path, arcname)
    finally:
        shutil.rmtree(temp_dir)

def modify_pptx_with_crop(pptx_path, left=0, top=0, right=0, bottom=0):
    """
    Modify a PPTX file to add crop to the first image.
    Crop values are in units of 100000.
    """
    temp_dir = tempfile.mkdtemp()
    try:
        # Extract PPTX
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Modify slide1.xml to add crop to first picture
        slide_path = os.path.join(temp_dir, 'ppt/slides/slide1.xml')
        if os.path.exists(slide_path):
            tree = etree.parse(slide_path)
            root = tree.getroot()
            
            # Find first picture
            pics = root.xpath('.//p:pic', namespaces=NSMAP)
            if pics:
                pic = pics[0]
                blipFill = pic.find('.//p:blipFill', NSMAP)
                if blipFill is not None:
                    srcRect = blipFill.find('.//a:srcRect', NSMAP)
                    if srcRect is None:
                        # Create srcRect if it doesn't exist
                        srcRect = etree.SubElement(blipFill, '{' + NSMAP['a'] + '}srcRect')
                    
                    # Set crop attributes
                    if left:
                        srcRect.set('l', str(left))
                    if top:
                        srcRect.set('t', str(top))
                    if right:
                        srcRect.set('r', str(right))
                    if bottom:
                        srcRect.set('b', str(bottom))
            
            # Write back
            tree.write(slide_path, xml_declaration=True, encoding='UTF-8', standalone=True)
        
        # Repack PPTX
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
            for root_dir, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_ref.write(file_path, arcname)
    finally:
        shutil.rmtree(temp_dir)

def create_base_presentation_with_image():
    """Create a presentation with one image on slide1."""
    prs = Presentation()
    
    # Add slide with blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add a test image (create a simple PNG)
    import tempfile
    from PIL import Image
    temp_img = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    img = Image.new('RGB', (100, 100), color='red')
    img.save(temp_img.name)
    temp_img.close()
    
    # Add image to slide
    left = top = Inches(1)
    pic = slide.shapes.add_picture(temp_img.name, left, top, width=Inches(2), height=Inches(2))
    
    # Clean up temp image
    os.unlink(temp_img.name)
    
    return prs

def create_geometry_fixtures():
    """Create geometry test fixtures."""
    fixtures_dir = Path("testdata/pptx/geometry")
    fixtures_dir.mkdir(parents=True, exist_ok=True)
    
    # Fixture 1: Image with 90-degree rotation
    print("Creating fixture: image with 90-degree rotation...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "rotation-90" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_rotation(str(fixture_path), 90)
    print(f"  Created: {fixture_path}")
    
    # Fixture 2: Image with 45-degree rotation
    print("Creating fixture: image with 45-degree rotation...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "rotation-45" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_rotation(str(fixture_path), 45)
    print(f"  Created: {fixture_path}")
    
    # Fixture 3: Image with horizontal flip
    print("Creating fixture: image with horizontal flip...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "flip-h" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_flip(str(fixture_path), flip_h=True)
    print(f"  Created: {fixture_path}")
    
    # Fixture 4: Image with vertical flip
    print("Creating fixture: image with vertical flip...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "flip-v" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_flip(str(fixture_path), flip_v=True)
    print(f"  Created: {fixture_path}")
    
    # Fixture 5: Image with both flips
    print("Creating fixture: image with both flips...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "flip-both" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_flip(str(fixture_path), flip_h=True, flip_v=True)
    print(f"  Created: {fixture_path}")
    
    # Fixture 6: Image with crop
    print("Creating fixture: image with crop...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "crop" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_crop(str(fixture_path), left=10000, top=20000, right=30000, bottom=40000)
    print(f"  Created: {fixture_path}")
    
    # Fixture 7: Image with rotation and flip
    print("Creating fixture: image with rotation and flip...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "rotation-and-flip" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_rotation(str(fixture_path), 90)
    modify_pptx_with_flip(str(fixture_path), flip_h=True)
    print(f"  Created: {fixture_path}")
    
    # Fixture 8: Image with all properties
    print("Creating fixture: image with rotation, flip, and crop...")
    prs = create_base_presentation_with_image()
    fixture_path = fixtures_dir / "all-properties" / "presentation.pptx"
    fixture_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(fixture_path))
    modify_pptx_with_rotation(str(fixture_path), 45)
    modify_pptx_with_flip(str(fixture_path), flip_h=True)
    modify_pptx_with_crop(str(fixture_path), left=5000, right=10000)
    print(f"  Created: {fixture_path}")

if __name__ == '__main__':
    try:
        create_geometry_fixtures()
        print("\nGeometry fixtures created successfully!")
    except Exception as e:
        print(f"Error creating fixtures: {e}", file=sys.stderr)
        sys.exit(1)
