"""
Generates a branded template PPTX for template capture and compilation testing.
This template includes:
- Title Slide archetype (title + subtitle)
- Content Slide archetype (title + body + optional image)
- Custom theme colors for brand consistency

Output: testdata/pptx/template-branded/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def add_title_slide(prs):
    """Add a branded title slide to the presentation."""
    # Use the title slide layout
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = "Title Slide"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)  # Brand blue
    
    # Set subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Subtitle"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)  # Brand red
    
    return slide


def add_content_slide(prs):
    """Add a branded content slide to the presentation."""
    # Use the title and content layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = "Content Slide"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 121)  # Brand blue
    
    # Set body content (bulleted list)
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Add bullet points
    p = text_frame.paragraphs[0]
    p.text = "First point"
    p.level = 0
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "Second point"
    p.level = 0
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "Third point"
    p.level = 0
    p.font.size = Pt(18)
    
    return slide


def generate_branded_template():
    """Generate a branded template presentation."""
    prs = Presentation()
    
    # Customize theme colors
    # Access the slide master to modify theme (optional, for future enhancement)
    
    # Add slides
    add_title_slide(prs)
    add_content_slide(prs)
    
    # Create output directory
    output_dir = "testdata/pptx/template-branded"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_branded_template()
