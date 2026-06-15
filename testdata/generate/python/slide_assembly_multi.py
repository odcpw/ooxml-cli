"""
Generates a PPTX with multiple slides using varied layouts for testing slide assembly operations
(delete, move, reorder, import, merge).

Output: testdata/pptx/slide-assembly-multi/presentation.pptx
Slides: 5 slides with different layouts and content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os


def generate_slide_assembly_multi():
    prs = Presentation()
    
    # Slide 1: Title Slide layout
    title_slide_layout = prs.slide_layouts[0]  # "Title Slide"
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Multi-Slide Assembly Test Deck"
    subtitle.text = "Testing delete, move, and reorder operations"
    
    # Slide 2: Title and Content layout
    content_layout = prs.slide_layouts[1]  # "Title and Content"
    slide2 = prs.slides.add_slide(content_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "First Content Slide"
    text_frame = content.text_frame
    text_frame.text = "This is the first content slide with bullet points:"
    p = text_frame.add_paragraph()
    p.text = "Point one"
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "Point two"
    p.level = 1
    
    # Slide 3: Section Header (layout 2 if available)
    if len(prs.slide_layouts) > 2:
        section_layout = prs.slide_layouts[2]
    else:
        section_layout = prs.slide_layouts[1]
    
    slide3 = prs.slides.add_slide(section_layout)
    title = slide3.shapes.title
    title.text = "Section Header Slide"
    
    # Slide 4: Title and Content layout (another one)
    slide4 = prs.slides.add_slide(content_layout)
    title = slide4.shapes.title
    content = slide4.placeholders[1]
    title.text = "Second Content Slide"
    text_frame = content.text_frame
    text_frame.text = "This is another content slide:"
    p = text_frame.add_paragraph()
    p.text = "Additional point one"
    p.level = 1
    p = text_frame.add_paragraph()
    p.text = "Additional point two"
    p.level = 1
    
    # Slide 5: Blank layout
    blank_layout = prs.slide_layouts[6]  # "Blank"
    slide5 = prs.slides.add_slide(blank_layout)
    
    # Add custom content to blank slide
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5)
    text_box = slide5.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = "Blank Layout Slide"
    p = text_frame.add_paragraph()
    p.text = "This slide uses the blank layout"
    p.level = 0
    
    # Create output directory
    output_dir = "testdata/pptx/slide-assembly-multi"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Layouts used: {prs.slide_layouts[0].name}, {prs.slide_layouts[1].name}, {prs.slide_layouts[2].name if len(prs.slide_layouts) > 2 else 'N/A'}, {prs.slide_layouts[6].name}")


if __name__ == "__main__":
    generate_slide_assembly_multi()
