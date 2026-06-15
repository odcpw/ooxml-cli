"""
Generates a PPTX with a slide containing speaker notes.
Output: testdata/pptx/notes-slide/presentation.pptx
"""

from pptx import Presentation
import os


def generate_notes_slide():
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Notes Slide"
    
    # Add content slide with notes
    content_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "Slide with Speaker Notes"
    
    text_frame = content.text_frame
    text_frame.text = "Main content on slide"
    
    # Add speaker notes
    notes_slide = slide2.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "These are the speaker notes for this slide.\nThey should not be visible to the audience."
    
    # Create output directory
    output_dir = "testdata/pptx/notes-slide"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_notes_slide()
