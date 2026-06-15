"""
Generates a PPTX with one normal slide chart and an embedded workbook.
Output: testdata/pptx/chart-simple/presentation.pptx
"""

import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def add_chart_slide(prs, slide_title, chart_name, chart_title, categories, values):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = slide_title

    data = CategoryChartData()
    data.categories = categories
    data.add_series("Revenue", values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1.5),
        Inches(8),
        Inches(4.5),
        data,
    )
    chart_shape.name = chart_name
    chart = chart_shape.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title


def generate_chart_simple():
    prs = Presentation()

    add_chart_slide(
        prs,
        "Chart Data",
        "Revenue Chart",
        "Revenue by Region",
        ["North", "South", "West"],
        (100, 120, 140),
    )
    add_chart_slide(
        prs,
        "Chart Data 2",
        "Revenue Chart Q2",
        "Revenue by Region Q2",
        ["North", "South", "West"],
        (130, 160, 190),
    )

    output_dir = "testdata/pptx/chart-simple"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"Generated {output_path}")


if __name__ == "__main__":
    generate_chart_simple()
