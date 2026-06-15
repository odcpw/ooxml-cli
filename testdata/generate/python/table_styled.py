"""
Generates a PPTX with a styled table (colors, borders, formatting).
Output: testdata/pptx/table-styled/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os


def generate_table_styled():
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Styled Table"
    
    # Add table slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add 3x3 table
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Style header row (row 0)
    for j in range(cols):
        cell = table_shape.cell(0, j)
        cell.text = f"Header {j}"
        
        # Set text formatting
        text_frame = cell.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Set cell background color
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue background
    
    # Fill remaining cells with data
    for i in range(1, rows):
        for j in range(cols):
            cell = table_shape.cell(i, j)
            cell.text = f"Data {i}-{j}"
            
            # Alternating row colors
            if i % 2 == 0:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(230, 240, 250)  # Light blue
            
            # Format text
            text_frame = cell.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    # Create output directory
    output_dir = "testdata/pptx/table-styled"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_table_styled()
