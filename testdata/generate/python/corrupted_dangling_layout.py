"""
Generates a corrupted PPTX where a slide references a non-existent slide layout.
This fixture tests the CLI's ability to detect broken layout references.
Output: testdata/pptx/corrupted-dangling-layout/presentation.pptx

Provenance: Python-pptx generated base + manual ZIP manipulation to break layout reference.
Corruption: prs/slides/slide2.xml references prs/slideLayouts/slideLayout99.xml which does not exist.
"""

from pptx import Presentation
import os
import zipfile
import shutil
from xml.etree import ElementTree as ET


def generate_corrupted_dangling_layout():
    # First create a normal PPTX
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Corrupted: Dangling Layout"
    
    # Add second slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    title.text = "This slide has broken layout reference"
    
    # Create temporary output directory
    output_dir = "testdata/pptx/corrupted-dangling-layout"
    os.makedirs(output_dir, exist_ok=True)
    
    temp_pptx = os.path.join(output_dir, "_temp.pptx")
    final_pptx = os.path.join(output_dir, "presentation.pptx")
    
    # Save the temporary PPTX
    prs.save(temp_pptx)
    
    # Now corrupt it by modifying the ZIP
    # Extract the PPTX (it's a ZIP file)
    extract_dir = os.path.join(output_dir, "_extracted")
    if os.path.exists(extract_dir):
        shutil.rmtree(extract_dir)
    os.makedirs(extract_dir)
    
    with zipfile.ZipFile(temp_pptx, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
    
    # Modify slide2.xml.rels to reference a non-existent layout
    slide2_rels_path = os.path.join(extract_dir, "ppt", "slides", "_rels", "slide2.xml.rels")
    if os.path.exists(slide2_rels_path):
        with open(slide2_rels_path, 'r', encoding='utf-8') as f:
            rels_content = f.read()
        
        # Replace the layout relationship with a reference to a non-existent layout
        # Change the slideLayout relationship to point to slideLayout99
        rels_content = rels_content.replace(
            'slideLayout',
            'slideLayout99'
        )
        
        with open(slide2_rels_path, 'w', encoding='utf-8') as f:
            f.write(rels_content)
    
    # Modify slide2.xml to have a comment about the broken layout
    slide2_xml_path = os.path.join(extract_dir, "ppt", "slides", "slide2.xml")
    if os.path.exists(slide2_xml_path):
        with open(slide2_xml_path, 'r', encoding='utf-8') as f:
            xml_content = f.read()
        
        # Add a comment noting the corruption
        xml_content = xml_content.replace(
            '<?xml',
            '<?xml'
        )  # Just a placeholder - the real corruption is in the .rels file
        
        with open(slide2_xml_path, 'w', encoding='utf-8') as f:
            f.write(xml_content)
    
    # Re-create the PPTX from the extracted (and modified) files
    with zipfile.ZipFile(final_pptx, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, dirs, files in os.walk(extract_dir):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, extract_dir)
                zipf.write(file_path, arcname)
    
    # Clean up temporary files
    os.remove(temp_pptx)
    shutil.rmtree(extract_dir)
    
    print(f"✓ Generated corrupted fixture: {final_pptx}")
    print(f"  Corruption: slide2.xml references slideLayout99 which does not exist")


if __name__ == "__main__":
    generate_corrupted_dangling_layout()
