"""
Generates a PPTX with a table containing merged cells.
Output: testdata/pptx/table-merged/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
import os


def generate_table_merged():
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Table with Merged Cells"
    
    # Add table slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add 4x4 table to demonstrate merges
    rows, cols = 4, 4
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3.5)
    
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Fill table with data
    for i in range(rows):
        for j in range(cols):
            cell = table_shape.cell(i, j)
            cell.text = f"R{i}C{j}"
    
    # Merge cells: merge cells (0,0) to (0,1) horizontally
    cell_a = table_shape.cell(0, 0)
    cell_b = table_shape.cell(0, 1)
    cell_a.merge(cell_b)
    cell_a.text = "Merged: R0C0-C1"
    
    # Merge cells: merge cells (1,2) to (2,2) vertically
    cell_c = table_shape.cell(1, 2)
    cell_d = table_shape.cell(2, 2)
    cell_c.merge(cell_d)
    cell_c.text = "Merged: R1-R2C2"
    
    # Create output directory
    output_dir = "testdata/pptx/table-merged"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_table_merged()
