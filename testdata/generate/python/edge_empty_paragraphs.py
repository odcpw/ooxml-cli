"""
Generates a PPTX with empty paragraphs (edge case for rich text handling).
Output: testdata/pptx/edge-empty-paragraphs/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Pt
import os


def generate_edge_empty_paragraphs():
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Empty Paragraphs Test"
    
    # Slide 2: Content with empty paragraphs
    content_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    title.text = "Content with Empty Lines"
    
    content = slide2.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Add paragraph with text
    p = text_frame.paragraphs[0]
    p.text = "First paragraph with text"
    
    # Add empty paragraph
    p = text_frame.add_paragraph()
    p.text = ""
    
    # Add another paragraph with text
    p = text_frame.add_paragraph()
    p.text = "Third paragraph (after empty)"
    
    # Add multiple empty paragraphs
    for _ in range(2):
        p = text_frame.add_paragraph()
        p.text = ""
    
    # Add final paragraph
    p = text_frame.add_paragraph()
    p.text = "Last paragraph"
    
    # Create output directory
    output_dir = "testdata/pptx/edge-empty-paragraphs"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_edge_empty_paragraphs()
