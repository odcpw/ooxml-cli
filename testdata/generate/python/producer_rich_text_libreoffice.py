"""
Generates a PPTX that would be created by LibreOffice with rich text.
Output: testdata/pptx/producers/libreoffice-rich-text/presentation.pptx

Note: Since python-pptx generates LibreOffice-compatible PPTX, this is
effectively the same as PowerPoint output but marked for LibreOffice testing.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def generate_libreoffice_rich_text():
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "LibreOffice Rich Text"
    subtitle = slide1.placeholders[1]
    subtitle.text = "Formatting Test"
    
    # Slide 2: Content with rich text
    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    title = slide2.shapes.title
    title.text = "Rich Text Examples"
    
    content = slide2.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Paragraph 1: Basic formatting
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.text = "Basic text with formatting"
    for run in p.runs:
        run.font.size = Pt(14)
    
    # Paragraph 2: Bold and colors
    p = text_frame.add_paragraph()
    p.level = 1
    
    run = p.add_run()
    run.text = "Bold and colored: "
    run.font.bold = True
    
    run = p.add_run()
    run.text = "Red text"
    run.font.color.rgb = RGBColor(255, 0, 0)
    
    # Paragraph 3: Different sizes
    p = text_frame.add_paragraph()
    p.level = 0
    p.text = "Multiple sizes: 10pt 12pt 14pt 16pt"
    
    sizes = [10, 12, 14, 16]
    text_parts = ["10pt ", "12pt ", "14pt ", "16pt"]
    
    p.clear()
    for size, text in zip(sizes, text_parts):
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
    
    # Paragraph 4: Italic and underline
    p = text_frame.add_paragraph()
    p.level = 1
    p.text = "Italic and underlined"
    
    for run in p.runs:
        run.font.italic = True
        run.font.underline = True
    
    # Paragraph 5: Mixed formatting
    p = text_frame.add_paragraph()
    p.level = 2
    
    run = p.add_run()
    run.text = "Mixed: "
    
    run = p.add_run()
    run.text = "Normal "
    
    run = p.add_run()
    run.text = "Bold "
    run.font.bold = True
    
    run = p.add_run()
    run.text = "Italic "
    run.font.italic = True
    
    run = p.add_run()
    run.text = "Colored"
    run.font.color.rgb = RGBColor(0, 128, 0)
    
    # Create output directory
    output_dir = "testdata/pptx/producers/libreoffice-rich-text"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_libreoffice_rich_text()
