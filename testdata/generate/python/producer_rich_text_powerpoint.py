"""
Generates a PPTX created by PowerPoint with rich text formatting.
Output: testdata/pptx/producers/powerpoint-rich-text/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os


def generate_powerpoint_rich_text():
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "PowerPoint Rich Text"
    subtitle = slide1.placeholders[1]
    subtitle.text = "Formatting Test"
    
    # Slide 2: Content with rich text
    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    title = slide2.shapes.title
    title.text = "Rich Text Examples"
    
    content = slide2.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Paragraph 1: Mixed formatting
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    
    run = p.add_run()
    run.text = "Normal text "
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    run = p.add_run()
    run.text = "Bold text "
    run.font.bold = True
    run.font.size = Pt(14)
    
    run = p.add_run()
    run.text = "Italic text "
    run.font.italic = True
    run.font.size = Pt(14)
    
    run = p.add_run()
    run.text = "Colored text"
    run.font.color.rgb = RGBColor(255, 0, 0)
    run.font.size = Pt(14)
    
    # Paragraph 2: Different level and alignment
    p = text_frame.add_paragraph()
    p.text = "Sub-level bullet with larger font"
    p.level = 1
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(16)
    
    # Paragraph 3: Custom colors and sizes
    p = text_frame.add_paragraph()
    p.level = 0
    p.alignment = PP_ALIGN.LEFT
    
    run = p.add_run()
    run.text = "Tiny "
    run.font.size = Pt(8)
    
    run = p.add_run()
    run.text = "Small "
    run.font.size = Pt(10)
    
    run = p.add_run()
    run.text = "Medium "
    run.font.size = Pt(12)
    
    run = p.add_run()
    run.text = "Large "
    run.font.size = Pt(18)
    
    run = p.add_run()
    run.text = "Huge"
    run.font.size = Pt(24)
    
    # Paragraph 4: Underline and strikethrough
    p = text_frame.add_paragraph()
    p.text = "Formatted text"
    p.level = 2
    
    for run in p.runs:
        run.font.underline = True
        run.font.color.rgb = RGBColor(0, 0, 255)
    
    # Create output directory
    output_dir = "testdata/pptx/producers/powerpoint-rich-text"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_powerpoint_rich_text()
