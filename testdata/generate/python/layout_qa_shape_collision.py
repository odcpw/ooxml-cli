"""
Generates a PPTX with a slide containing overlapping shapes.
This fixture is used to test shape collision detection heuristics.
Output: testdata/pptx/layout-qa-shape-collision/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import os


def generate_shape_collision():
    prs = Presentation()
    
    # Use a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add first shape (rectangle)
    left1 = Inches(1)
    top1 = Inches(1)
    width1 = Inches(2)
    height1 = Inches(2)
    
    shape1 = slide.shapes.add_shape(
        1,  # Rectangle shape type
        left1, top1, width1, height1
    )
    shape1.name = "Rectangle 1"
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
    shape1.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    
    # Add second shape that overlaps with the first
    # This overlap should be detected
    left2 = Inches(1.5)  # Overlaps with first shape
    top2 = Inches(1.5)   # Overlaps with first shape
    width2 = Inches(2)
    height2 = Inches(2)
    
    shape2 = slide.shapes.add_shape(
        1,  # Rectangle shape type
        left2, top2, width2, height2
    )
    shape2.name = "Rectangle 2"
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue
    shape2.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    
    # Add a third shape with significant overlap
    left3 = Inches(1.2)
    top3 = Inches(1.2)
    width3 = Inches(2)
    height3 = Inches(2)
    
    shape3 = slide.shapes.add_shape(
        1,  # Rectangle shape type
        left3, top3, width3, height3
    )
    shape3.name = "Rectangle 3"
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(0, 255, 0)  # Green
    shape3.line.color.rgb = RGBColor(0, 0, 0)  # Black border
    
    # Create output directory
    output_dir = "testdata/pptx/layout-qa-shape-collision"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_shape_collision()
