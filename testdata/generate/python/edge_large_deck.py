"""
Generates a PPTX with many slides (large deck edge case).
Output: testdata/pptx/edge-large-deck/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os


def generate_edge_large_deck():
    prs = Presentation()
    
    # Add 50 slides with varied content
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Large Deck Test"
    subtitle = slide1.placeholders[1]
    subtitle.text = "50 slides"
    
    # Add 49 more content slides
    for i in range(2, 51):
        slide = prs.slides.add_slide(content_layout)
        title = slide.shapes.title
        title.text = f"Slide {i}"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = f"Content for slide {i}\n\n"
        
        # Add some text content
        for j in range(3):
            p = text_frame.add_paragraph()
            p.text = f"Bullet point {j+1} on slide {i}"
            p.level = 0 if j % 2 == 0 else 1
    
    # Create output directory
    output_dir = "testdata/pptx/edge-large-deck"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path} (50 slides)")


if __name__ == "__main__":
    generate_edge_large_deck()
