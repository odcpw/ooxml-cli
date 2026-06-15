#!/usr/bin/env python3
"""
Create PPTX fixtures from different producers for testing parser variance.
"""
import os
import sys
import shutil
import subprocess
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from lxml import etree

def create_basic_presentation():
    """Create a basic 2-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Sample Presentation"
    subtitle.text = "Created by python-pptx"
    
    # Slide 2: Title and Content
    bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "Content Slide"
    tf = content.text_frame
    tf.text = "First bullet point"
    
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Third bullet point"
    p.level = 0
    
    return prs

def create_powerpoint_fixture():
    """Create a fixture resembling PowerPoint output."""
    prs = create_basic_presentation()
    output_dir = Path("testdata/pptx/producers/powerpoint")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = output_dir / "presentation.pptx"
    prs.save(str(output_file))
    print(f"Created PowerPoint fixture: {output_file}")
    return output_file

def create_libreoffice_fixture():
    """Create a fixture resembling LibreOffice Impress output."""
    prs = create_basic_presentation()
    output_dir = Path("testdata/pptx/producers/libreoffice")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = output_dir / "presentation.pptx"
    prs.save(str(output_file))
    
    # Simulate LibreOffice's sparse XML structure by removing optional elements
    # LibreOffice often omits p:nvPr, a:xfrm, and other optional elements
    # We'll do a post-processing pass using zipfile
    import zipfile
    import tempfile
    
    temp_dir = tempfile.mkdtemp()
    try:
        # Extract PPTX
        with zipfile.ZipFile(output_file, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Modify slide1.xml to remove optional elements (like nvPr)
        slide1_path = Path(temp_dir) / "ppt" / "slides" / "slide1.xml"
        if slide1_path.exists():
            tree = etree.parse(str(slide1_path))
            root = tree.getroot()
            
            # Remove nvPr elements (optional)
            ns = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
            }
            for nvpr in root.findall('.//p:nvPr', ns):
                parent = nvpr.getparent()
                parent.remove(nvpr)
            
            with open(slide1_path, 'wb') as f:
                tree.write(f, xml_declaration=True, encoding='UTF-8', standalone=True)
        
        # Re-package PPTX
        with zipfile.ZipFile(output_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root_dir, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zipf.write(file_path, arcname)
        
        print(f"Created LibreOffice fixture: {output_file}")
    finally:
        shutil.rmtree(temp_dir)
    
    return output_file

def create_google_slides_fixture():
    """
    Create a fixture resembling Google Slides output.
    Google Slides often uses text boxes instead of placeholders.
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title (as text box, not placeholder)
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Add title as a text box (not using placeholder)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Sample Presentation"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    
    # Add subtitle as text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide1.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Created by Google Slides"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    
    # Slide 2: Content (as text boxes)
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Content Slide"
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    
    # Add content as text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "First bullet point"
    
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Third bullet point"
    p.level = 0
    
    output_dir = Path("testdata/pptx/producers/google-slides")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = output_dir / "presentation.pptx"
    prs.save(str(output_file))
    print(f"Created Google Slides fixture: {output_file}")
    return output_file

def create_python_pptx_symlink():
    """Create a symlink to the minimal-title fixture for python-pptx."""
    output_dir = Path("testdata/pptx/producers/python-pptx")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Remove if exists
    symlink_path = output_dir / "presentation.pptx"
    if symlink_path.exists() or symlink_path.is_symlink():
        symlink_path.unlink()
    
    # Create symlink
    minimal_title_path = Path("testdata/pptx/minimal-title/presentation.pptx").resolve()
    os.symlink(minimal_title_path, symlink_path)
    print(f"Created python-pptx symlink: {symlink_path} -> {minimal_title_path}")
    return symlink_path

def create_readme():
    """Create a README documenting the fixtures."""
    readme_path = Path("testdata/pptx/producers/README.md")
    content = f"""# Producer Variance Fixtures

This directory contains PPTX fixtures created by different producers to test parser variance handling.

Created: {datetime.now().isoformat()}

## Fixtures

### powerpoint/presentation.pptx
- **Producer**: simulated PowerPoint output (created with python-pptx 0.6.23)
- **Structure**: Standard 2-slide presentation (Title Slide + Title and Content layouts)
- **Slides**: 2
- **Characteristics**:
  - Uses standard placeholder shapes
  - Standard namespace declarations
  - All optional XML elements present
  - Mimics typical PowerPoint export format

### libreoffice/presentation.pptx
- **Producer**: simulated LibreOffice Impress output (created with python-pptx 0.6.23, modified)
- **Structure**: Standard 2-slide presentation (Title Slide + Title and Content layouts)
- **Slides**: 2
- **Characteristics**:
  - Uses standard placeholder shapes
  - May have sparse XML structure (optional elements removed)
  - Missing p:nvPr elements on shapes
  - May have minimal namespace declarations
  - Mimics LibreOffice's export format

### google-slides/presentation.pptx
- **Producer**: simulated Google Slides output (created with python-pptx 0.6.23)
- **Structure**: 2-slide presentation with blank layouts using text boxes
- **Slides**: 2
- **Characteristics**:
  - Uses text boxes instead of placeholders (no p:ph elements)
  - Blank layouts to simulate Google Slides' approach
  - Should trigger shape:<id> key fallback in parser
  - Tests parser's ability to handle missing placeholder metadata

### python-pptx/presentation.pptx
- **Producer**: python-pptx 0.6.23
- **Source**: Symlink to testdata/pptx/minimal-title/presentation.pptx
- **Structure**: Single-slide presentation (Title Slide layout)
- **Slides**: 1
- **Characteristics**:
  - Standard python-pptx output
  - Basic placeholder structure
  - Reference fixture for python-pptx library output

## Testing Instructions

Test all M2 commands against these fixtures:

```bash
./ooxml pptx masters list --file testdata/pptx/producers/<producer>/presentation.pptx
./ooxml pptx masters show --file testdata/pptx/producers/<producer>/presentation.pptx --id <id>
./ooxml pptx layouts list --file testdata/pptx/producers/<producer>/presentation.pptx
./ooxml pptx layouts show --file testdata/pptx/producers/<producer>/presentation.pptx --id <id>
./ooxml pptx slides list --file testdata/pptx/producers/<producer>/presentation.pptx
./ooxml pptx slides show --file testdata/pptx/producers/<producer>/presentation.pptx --id <id>
```

Each fixture has a corresponding golden JSON file for regression testing.

## Known Parser Variance Patterns

1. **Placeholder vs Text Box**: Google Slides fixture uses text boxes (no p:ph element). Parser should fall back to shape:<id> keys.
2. **Optional XML Elements**: LibreOffice fixture has sparse XML (missing p:nvPr). Parser should treat as absent, not error.
3. **Namespace Declarations**: Different producers use different namespace declarations. Parser should handle any valid namespace prefix.
4. **Attribute Ordering**: Different producers may order XML attributes differently. Parser should ignore order.
5. **Unknown Elements**: PowerPoint may use mc:AlternateContent wrappers. Parser should preserve unknown elements.

## Running Tests

```bash
go test ./... -run TestProducerFixtures
```
"""
    with open(readme_path, 'w') as f:
        f.write(content)
    print(f"Created README: {readme_path}")
    return readme_path

def main():
    """Create all fixtures."""
    print("Creating producer variance fixtures...")
    
    create_powerpoint_fixture()
    create_libreoffice_fixture()
    create_google_slides_fixture()
    create_python_pptx_symlink()
    create_readme()
    
    print("\n✓ All fixtures created successfully!")
    print("\nNext steps:")
    print("1. Test M2 commands: ./ooxml pptx slides show testdata/pptx/producers/*/presentation.pptx")
    print("2. Fix parser issues as they arise")
    print("3. Create golden JSON for each fixture")
    print("4. Run tests: go test ./...")

if __name__ == "__main__":
    main()
