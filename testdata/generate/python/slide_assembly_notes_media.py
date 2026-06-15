"""
Generates a PPTX with slides containing speaker notes and embedded media for slide assembly testing.

Output: testdata/pptx/slide-assembly-notes-media/presentation.pptx
Purpose: Test fixture for slides with notes and embedded images during assembly operations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os
import tempfile


def create_test_image(width=400, height=300, color='red', label='Image'):
    """Create a simple test image and return path."""
    img = Image.new('RGB', (width, height), color=color)
    # Add some variation with a rectangle
    pixels = img.load()
    for x in range(50, 150):
        for y in range(50, 150):
            pixels[x, y] = (255 - ord(label[0]) % 256, 200, 100)
    
    # Save to a temporary location
    img_path = os.path.join(tempfile.gettempdir(), f"test_image_{label.replace(' ', '_')}.png")
    img.save(img_path)
    return img_path


def generate_slide_assembly_notes_media():
    prs = Presentation()
    
    # Create test images
    image1_path = create_test_image(400, 300, 'red', 'Image1')
    image2_path = create_test_image(400, 300, 'green', 'Image2')
    image3_path = create_test_image(400, 300, 'blue', 'Image3')
    
    try:
        # Slide 1: Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "Notes and Media Test"
        subtitle.text = "Slides with notes and embedded images"
        
        # Add notes to title slide
        notes_slide = slide1.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "This is the opening slide.\nPresenter notes are included for reference."
        
        # Slide 2: Content with image and notes
        content_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(content_layout)
        title = slide2.shapes.title
        title.text = "Slide with Image and Notes"
        
        # Add embedded image
        left = Inches(4.5)
        top = Inches(1.5)
        width = Inches(4)
        height = Inches(3)
        slide2.shapes.add_picture(image1_path, left, top, width, height)
        
        # Add text content
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(3.5)
        height = Inches(4)
        text_box = slide2.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Content with embedded image:"
        p = text_frame.add_paragraph()
        p.text = "First image included in this slide"
        p.level = 1
        
        # Add detailed notes
        notes_slide = slide2.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "This slide demonstrates media inclusion.\nThe image shows test content.\nUse this for testing media preservation during assembly operations."
        
        # Slide 3: Multiple images with notes
        slide3 = prs.slides.add_slide(content_layout)
        title = slide3.shapes.title
        title.text = "Multiple Images and Notes"
        
        # Add first image
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(3)
        height = Inches(2.5)
        slide3.shapes.add_picture(image2_path, left, top, width, height)
        
        # Add second image
        left = Inches(3.8)
        top = Inches(1.5)
        width = Inches(3)
        height = Inches(2.5)
        slide3.shapes.add_picture(image3_path, left, top, width, height)
        
        # Add caption text
        left = Inches(0.5)
        top = Inches(4.2)
        width = Inches(6.3)
        height = Inches(1)
        text_box = slide3.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Two images demonstrating multi-media slide content"
        
        # Add notes
        notes_slide = slide3.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "This slide includes multiple images.\nBoth images should be preserved when slides are moved or imported.\nThis tests media relationship handling in assembly operations."
        
        # Slide 4: Content with mixed notes
        slide4 = prs.slides.add_slide(content_layout)
        title = slide4.shapes.title
        title.text = "Complex Notes Example"
        
        # Add content placeholder
        content = slide4.placeholders[1]
        text_frame = content.text_frame
        text_frame.text = "Slide with comprehensive speaker notes"
        p = text_frame.add_paragraph()
        p.text = "Multiple paragraphs in content"
        p.level = 1
        
        # Add embedded image
        left = Inches(4.5)
        top = Inches(2)
        width = Inches(3.5)
        height = Inches(2.8)
        slide4.shapes.add_picture(image1_path, left, top, width, height)
        
        # Add comprehensive notes
        notes_slide = slide4.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "Complex notes structure for testing.\nFirst paragraph of notes.\nSecond paragraph explaining the content.\nThird paragraph with additional context.\nFinal paragraph testing multi-line note preservation."
        
        # Slide 5: Minimal slide with notes
        blank_layout = prs.slide_layouts[6]
        slide5 = prs.slides.add_slide(blank_layout)
        
        # Add text
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(5)
        text_box = slide5.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Final Slide"
        p = text_frame.add_paragraph()
        p.text = "This slide has simple notes"
        
        # Add minimal notes
        notes_slide = slide5.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "Simple closing notes for the final slide."
        
        # Create output directory
        output_dir = "testdata/pptx/slide-assembly-notes-media"
        os.makedirs(output_dir, exist_ok=True)
        
        # Save presentation
        output_path = os.path.join(output_dir, "presentation.pptx")
        prs.save(output_path)
        print(f"✓ Generated {output_path}")
        print(f"  Slides: {len(prs.slides)}")
        print(f"  Features: Notes on all slides, embedded images on slides 2-4")
        
    finally:
        # Clean up temporary image files
        for img_path in [image1_path, image2_path, image3_path]:
            if os.path.exists(img_path):
                try:
                    os.remove(img_path)
                except OSError:
                    pass


if __name__ == "__main__":
    generate_slide_assembly_notes_media()
