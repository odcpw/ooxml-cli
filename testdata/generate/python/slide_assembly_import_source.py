"""
Generates a PPTX with slides using a different theme and layout set for import/merge testing.

Output: testdata/pptx/slide-assembly-import-source/presentation.pptx
Purpose: Source deck for testing cross-deck slide imports with different themes and layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import os


def generate_slide_assembly_import_source():
    prs = Presentation()
    
    # Slide 1: Title Slide with custom theme color
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Import Source Deck"
    subtitle.text = "With custom theme and styling for import tests"
    
    # Change title color to demonstrate theme difference
    title_frame = title.text_frame
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 102, 204)  # Blue
    
    # Slide 2: Title and Content with custom layout
    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    title.text = "Import Source Content 1"
    text_frame = content.text_frame
    text_frame.text = "This slide is from a different deck:"
    
    # Add colored paragraph
    p = text_frame.add_paragraph()
    p.text = "Content with custom styling"
    p.level = 1
    for run in p.runs:
        run.font.color.rgb = RGBColor(204, 0, 0)  # Red
        run.font.bold = True
    
    # Slide 3: Another layout variation
    if len(prs.slide_layouts) > 3:
        layout3 = prs.slide_layouts[3]
    else:
        layout3 = prs.slide_layouts[1]
    
    slide3 = prs.slides.add_slide(layout3)
    title = slide3.shapes.title
    title.text = "Import Source Content 2"
    
    # Add a shape with custom color
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(5)
    height = Inches(3)
    text_box = slide3.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = "This content will be imported with different styling"
    
    # Apply fill color to the text box
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(220, 220, 220)  # Light gray
    
    # Slide 4: Blank layout for variety
    blank_layout = prs.slide_layouts[6]
    slide4 = prs.slides.add_slide(blank_layout)
    
    # Add custom shapes with different colors
    left = Inches(1)
    top = Inches(1)
    width = Inches(3)
    height = Inches(4)
    text_box = slide4.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "Custom Layout Slide from Import Source"
    
    # Create output directory
    output_dir = "testdata/pptx/slide-assembly-import-source"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Layouts: Different theme with custom colors and styling")


if __name__ == "__main__":
    generate_slide_assembly_import_source()
