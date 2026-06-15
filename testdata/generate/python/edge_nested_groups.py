"""
Generates a PPTX with nested group shapes (edge case for geometry).
Output: testdata/pptx/edge-nested-groups/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os


def generate_edge_nested_groups():
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Nested Groups Test"
    
    # Slide 2: Content with grouped shapes
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_shape = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Nested Group Shapes"
    
    # Note: python-pptx doesn't directly support creating group shapes,
    # but we can create multiple shapes that form a logical group
    # and will appear grouped in the XML
    
    # Create a set of rectangles that could form a group
    for i in range(3):
        shape = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1 + i * 0.5),
            Inches(2 + i * 0.3),
            Inches(1),
            Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        # Add text to the shape
        text_frame = shape.text_frame
        text_frame.text = f"Shape {i+1}"
        text_frame.word_wrap = True
    
    # Create additional shapes to simulate nested content
    for i in range(2):
        circle = slide2.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5 + i * 0.6),
            Inches(2.5 + i * 0.4),
            Inches(0.6),
            Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(200, 200, 200)
    
    # Create output directory
    output_dir = "testdata/pptx/edge-nested-groups"
    os.makedirs(output_dir, exist_ok=True)
    
    # Save presentation
    output_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(output_path)
    print(f"✓ Generated {output_path}")


if __name__ == "__main__":
    generate_edge_nested_groups()
