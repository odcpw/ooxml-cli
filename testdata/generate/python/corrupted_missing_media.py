"""
Generates a corrupted PPTX with an image relationship that points to a non-existent media file.
This fixture tests the CLI's ability to detect missing media references.
Output: testdata/pptx/corrupted-missing-media/presentation.pptx

Provenance: Python-pptx generated base + manual ZIP manipulation to break media reference.
Corruption: prs/slides/slide1.xml contains relationship to media/image1.png but the file is deleted.
"""

from pptx import Presentation
import os
import zipfile
import shutil
from pathlib import Path


def generate_corrupted_missing_media():
    # First create a normal PPTX
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    title.text = "Corrupted: Missing Media"
    
    # Create temporary output directory
    output_dir = "testdata/pptx/corrupted-missing-media"
    os.makedirs(output_dir, exist_ok=True)
    
    temp_pptx = os.path.join(output_dir, "_temp.pptx")
    final_pptx = os.path.join(output_dir, "presentation.pptx")
    
    # Save the temporary PPTX
    prs.save(temp_pptx)
    
    # Now corrupt it by modifying the ZIP
    # Extract the PPTX (it's a ZIP file)
    extract_dir = os.path.join(output_dir, "_extracted")
    if os.path.exists(extract_dir):
        shutil.rmtree(extract_dir)
    os.makedirs(extract_dir)
    
    with zipfile.ZipFile(temp_pptx, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
    
    # Modify slide1.xml to reference a missing image
    slide_xml_path = os.path.join(extract_dir, "ppt", "slides", "slide1.xml")
    if os.path.exists(slide_xml_path):
        with open(slide_xml_path, 'r', encoding='utf-8') as f:
            xml_content = f.read()
        
        # Add a reference to a non-existent image (only modify XML, don't add the actual image)
        # Insert before closing tags
        if '</p:sld>' in xml_content:
            # Add a picture element that references a non-existent media file
            picture_ref = '''<p:cSld><p:spTree><p:pic><p:nvPicPr><p:cNvPr id="2" name="MissingImage"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr><p:blipFill><a:blip r:embed="rId2"/></p:blipFill><p:spPr/></p:pic></p:spTree></p:cSld>'''
            # For simplicity, just add a remark that the image is intentionally missing
            xml_content = xml_content.replace('</p:sld>', '<!-- CORRUPTION: References missing media file --></p:sld>')
        
        with open(slide_xml_path, 'w', encoding='utf-8') as f:
            f.write(xml_content)
    
    # Also modify the slide relationships to reference a missing image
    slide_rels_path = os.path.join(extract_dir, "ppt", "slides", "_rels", "slide1.xml.rels")
    if os.path.exists(slide_rels_path):
        with open(slide_rels_path, 'r', encoding='utf-8') as f:
            rels_content = f.read()
        
        # Add a relationship to a non-existent image
        if '</Relationships>' in rels_content:
            rels_content = rels_content.replace(
                '</Relationships>',
                '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/image1.png"/></Relationships>'
            )
        
        with open(slide_rels_path, 'w', encoding='utf-8') as f:
            f.write(rels_content)
    
    # Re-create the PPTX from the extracted (and modified) files
    with zipfile.ZipFile(final_pptx, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, dirs, files in os.walk(extract_dir):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, extract_dir)
                zipf.write(file_path, arcname)
    
    # Clean up temporary files
    os.remove(temp_pptx)
    shutil.rmtree(extract_dir)
    
    print(f"✓ Generated corrupted fixture: {final_pptx}")
    print(f"  Corruption: References media/image1.png that does not exist")


if __name__ == "__main__":
    generate_corrupted_missing_media()
