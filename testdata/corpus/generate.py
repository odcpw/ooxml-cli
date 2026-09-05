#!/usr/bin/env python3
"""Regenerate self-authored corpus; dependencies are creation-only, never CI dependencies.

python-pptx==1.0.2 python-docx==1.2.0 XlsxWriter==3.2.5 and LibreOffice.
All prose and data below are original test content, distributed under the repo license.
"""
from pathlib import Path
import subprocess
import tempfile
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import xlsxwriter

ROOT = Path(__file__).resolve().parent
for producer in ('python-pptx', 'python-docx', 'xlsxwriter', 'libreoffice'):
    (ROOT / producer).mkdir(exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Corpus quarterly review'
frame = slide.placeholders[1].text_frame
frame.text = 'Revenue increased'
frame.add_paragraph().text = 'Unicode: café — Zürich'
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Regional results'
table = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(7), Inches(2)).table
for row, values in enumerate([['Region', 'Revenue'], ['North', '120'], ['South', '95']]):
    for col, value in enumerate(values):
        table.cell(row, col).text = value
prs.save(ROOT / 'python-pptx/review.pptx')

doc = Document()
doc.add_heading('Corpus quarterly review', 0)
doc.add_heading('Results', 1)
p = doc.add_paragraph('Revenue ')
p.add_run('increased').bold = True
p.add_run(' in Zürich — café sales improved.')
doc.add_paragraph('Review assumptions', style='List Bullet')
table = doc.add_table(rows=1, cols=2)
table.style = 'Table Grid'
for cell, value in zip(table.rows[0].cells, ['Region', 'Revenue']):
    cell.text = value
for values in [('North', '120'), ('South', '95')]:
    for cell, value in zip(table.add_row().cells, values):
        cell.text = value
doc.sections[0].header.paragraphs[0].text = 'Corpus report'
doc.save(ROOT / 'python-docx/report.docx')

with xlsxwriter.Workbook(ROOT / 'xlsxwriter/sales.xlsx') as book:
    book.set_properties({'created': datetime(2026, 9, 5), 'author': 'ooxml-cli corpus'})
    sheet = book.add_worksheet('Sales')
    sheet.add_table('A1:C4', {'columns': [{'header': 'Region'}, {'header': 'Revenue'}, {'header': 'Cost'}],
                            'data': [['North', 120, 80], ['South', 95, 60], ['Zürich', 110, 75]]})
    sheet.write('A6', 'Total')
    sheet.write_formula('B6', '=SUM(B2:B4)', None, 325)
    sheet.freeze_panes(1, 0)
    sheet.set_column('A:C', 18)
    chart = book.add_chart({'type': 'column'})
    chart.add_series({'name': 'Revenue', 'categories': '=Sales!$A$2:$A$4', 'values': '=Sales!$B$2:$B$4'})
    sheet.insert_chart('E2', chart)

with tempfile.TemporaryDirectory(prefix='azureowl-lo-') as profile:
    for source in ['python-pptx/review.pptx', 'python-docx/report.docx', 'xlsxwriter/sales.xlsx']:
        path = ROOT / source
        subprocess.run(['soffice', '-env:UserInstallation=' + Path(profile).as_uri(), '--headless',
                        '--convert-to', path.suffix[1:], '--outdir', str(ROOT / 'libreoffice'), str(path)],
                       check=True, timeout=60)
